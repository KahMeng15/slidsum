#!/usr/bin/env python3
"""
SlidSum - Lecture Slides Content Extractor & AI Analyzer
by kahmeng | kahmeng15.github.io
"""
import os
import sys
import glob
import base64
import subprocess
import logging
from pathlib import Path
import time
import importlib.util
import re
import warnings
import atexit
import shutil
import shlex
import traceback
import io

warnings.filterwarnings("ignore")

logging.basicConfig(level=logging.ERROR, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

project_root = os.path.dirname(os.path.abspath(__file__))
temp_dir = os.path.join(project_root, "temp")
os.makedirs(temp_dir, exist_ok=True)

def cleanup_temp_dir():
    if os.path.exists(temp_dir):
        for filename in os.listdir(temp_dir):
            file_path = os.path.join(temp_dir, filename)
            try:
                if os.path.isfile(file_path) or os.path.islink(file_path):
                    os.unlink(file_path)
                elif os.path.isdir(file_path):
                    shutil.rmtree(file_path)
            except Exception:
                pass

atexit.register(cleanup_temp_dir)

# --- DEPENDENCY CHECK ---
def get_missing_deps():
    deps = {
        "dotenv": "python-dotenv",
        "openai": "openai",
        "google.genai": "google-genai",
        "rich": "rich",
        "inquirer": "inquirer",
        "fitz": "PyMuPDF",
        "pptx": "python-pptx",
        "PIL": "Pillow",
    }
    missing = []
    for module, pkg in deps.items():
        try:
            if importlib.util.find_spec(module.split('.')[0]) is None:
                missing.append(pkg)
        except (ModuleNotFoundError, AttributeError):
            missing.append(pkg)
    return missing

def install_dependencies(missing):
    try:
        from rich.console import Console
        from rich.panel import Panel
        from rich.progress import Progress, SpinnerColumn, TextColumn
        console = Console()
    except ImportError:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "rich", "inquirer"])
        from rich.console import Console
        from rich.panel import Panel
        from rich.progress import Progress, SpinnerColumn, TextColumn
        console = Console()

    console.print(Panel("[bold magenta1]SlidSum - Environment Setup[/bold magenta1]\n[dim]Installing required libraries...[/dim]", border_style="magenta3"))

    with Progress(SpinnerColumn(), TextColumn("[progress.description]{task.description}"), transient=True) as progress:
        task = progress.add_task(description="Installing missing libraries...", total=None)
        try:
            process = subprocess.Popen(
                [sys.executable, "-m", "pip", "install", "-r", "requirements.txt"],
                stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True
            )
            for line in process.stdout:
                if line.strip():
                    short = line.strip()[:60] + "..." if len(line.strip()) > 60 else line.strip()
                    progress.update(task, description=f"[bold gold3]Installing:[/bold gold3] [dim]{short}[/dim]")
            process.wait()
            if process.returncode != 0:
                raise Exception("pip install failed")
            progress.update(task, description="[bold spring_green3]✓ Environment ready![/bold spring_green3]")
            time.sleep(1)
        except Exception as e:
            console.print(f"[bold red]Error during setup:[/bold red] {e}")
            sys.exit(1)

    console.print("[bold spring_green3]Setup Complete![/bold spring_green3] Please restart the script.\n")
    sys.exit(0)

missing = get_missing_deps()
if missing:
    install_dependencies(missing)

# --- MAIN IMPORTS ---
try:
    from dotenv import load_dotenv
    from openai import OpenAI
    from google import genai
    from google.genai import types as genai_types
    from rich.console import Console
    from rich.panel import Panel
    from rich.status import Status
    from rich.table import Table
    import inquirer
    try:
        from inquirer import Separator
    except ImportError:
        Separator = None
    import fitz  # PyMuPDF
    from pptx import Presentation
    from PIL import Image
except (ImportError, Exception) as e:
    print(f"[-] Critical Error: {e}")
    print("[*] Try: pip install --force-reinstall -r requirements.txt")
    sys.exit(1)

console = Console()

# --- IDLE TIMEOUT WRAPPER ---
original_prompt = inquirer.prompt

def prompt_with_timeout(*args, **kwargs):
    import signal
    def timeout_handler(signum, frame):
        raise TimeoutError("Idle timeout reached.")
    signal.signal(signal.SIGALRM, timeout_handler)
    try:
        timeout = int(os.getenv("IDLE_TIMEOUT", "300"))
    except ValueError:
        timeout = 300
    if timeout > 0:
        signal.alarm(timeout)
    try:
        # Use plain stdout/stderr for inquirer to avoid Rich console interference
        old_stdout = sys.stdout
        old_stderr = sys.stderr
        sys.stdout = sys.__stdout__
        sys.stderr = sys.__stderr__
        try:
            return original_prompt(*args, **kwargs)
        finally:
            sys.stdout = old_stdout
            sys.stderr = old_stderr
    except TimeoutError:
        console.print("\n[bold gold3]Idle timeout reached. Auto-quitting...[/bold gold3]")
        sys.exit(0)
    finally:
        if timeout > 0:
            signal.alarm(0)

inquirer.prompt = prompt_with_timeout

# --- ENVIRONMENT SETUP ---
def setup_environment():
    load_dotenv()
    config = {
        "INPUT_DIR":           os.getenv("INPUT_DIR", "./input_slides"),
        "OUTPUT_DIR":          os.getenv("OUTPUT_DIR", "./output_files"),
        "PROMPTS_DIR":         os.getenv("PROMPTS_DIR", "./prompts"),
        "TEMP_DIR":            temp_dir,
        "FILE_PATTERN":        os.getenv("FILE_PATTERN", "*.pdf,*.pptx,*.ppt,*.png,*.jpg,*.jpeg"),
        "OPENAI_API_KEY":      os.getenv("OPENAI_API_KEY"),
        "GEMINI_API_KEY":      os.getenv("GEMINI_API_KEY"),
        "GEMINI_MODELS":       os.getenv("GEMINI_MODELS", "gemini-2.0-flash").split(","),
        "OPENAI_MODELS":       os.getenv("OPENAI_MODELS", "gpt-4o-mini").split(","),
        "SUMMARIZER_PROVIDER": os.getenv("SUMMARIZER_PROVIDER", "gemini").lower(),
        "MAX_IMAGE_DIM":       int(os.getenv("MAX_IMAGE_DIM", "1600")),
    }
    config["GEMINI_MODELS"] = [m.strip() for m in config["GEMINI_MODELS"] if m.strip()]
    config["OPENAI_MODELS"] = [m.strip() for m in config["OPENAI_MODELS"] if m.strip()]

    if config["SUMMARIZER_PROVIDER"] == "openai":
        if not config["OPENAI_API_KEY"] or config["OPENAI_API_KEY"] == "your_openai_api_key_here":
            console.print("[bold red]Warning:[/bold red] OPENAI_API_KEY not set in .env")
    elif config["SUMMARIZER_PROVIDER"] == "gemini":
        if not config["GEMINI_API_KEY"] or config["GEMINI_API_KEY"] == "your_gemini_api_key_here":
            console.print("[bold red]Warning:[/bold red] GEMINI_API_KEY not set in .env")

    os.makedirs(config["INPUT_DIR"], exist_ok=True)
    os.makedirs(config["OUTPUT_DIR"], exist_ok=True)
    os.makedirs(config["PROMPTS_DIR"], exist_ok=True)
    return config

# --- IMAGE UTILITIES ---
def resize_image(img_bytes: bytes, max_dim: int) -> bytes:
    img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
    w, h = img.size
    if max(w, h) > max_dim:
        ratio = max_dim / max(w, h)
        img = img.resize((int(w * ratio), int(h * ratio)), Image.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()

# --- SLIDE EXTRACTION ---
def extract_slides_content(file_path: str, config: dict):
    """
    Extract text and image parts from a slide file.
    Returns:
        content_text (str): Markdown text extracted slide-by-slide.
        image_parts (list): [(bytes, mime_type, label), ...] for image-based pages.
    """
    ext = Path(file_path).suffix.lower()
    max_dim = config.get("MAX_IMAGE_DIM", 1600)
    lines = []
    image_parts = []

    if ext == ".pdf":
        doc = fitz.open(file_path)
        total = len(doc)
        console.print(f"[dim]  PDF: {total} page(s) detected.[/dim]")
        for i, page in enumerate(doc, 1):
            text = page.get_text("text").strip()
            lines.append(f"## Page {i}")
            if text:
                lines.append(text)
            else:
                console.print(f"[dim]  Page {i}: no text — rendering as image for AI vision.[/dim]")
                pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
                img_bytes = resize_image(pix.tobytes("png"), max_dim)
                image_parts.append((img_bytes, "image/png", f"Page {i}"))
                lines.append(f"[Image-only page — attached as Page {i}]")
            lines.append("")
        doc.close()

    elif ext in (".pptx", ".ppt"):
        prs = Presentation(file_path)
        total = len(prs.slides)
        console.print(f"[dim]  PPTX: {total} slide(s) detected.[/dim]")
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"## Slide {i}")
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = " ".join(run.text for run in para.runs).strip()
                        if t:
                            slide_texts.append(t)
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_texts.append(f"\n**Speaker Notes:** {notes}")
            if slide_texts:
                lines.extend(slide_texts)
            else:
                lines.append("[No text content on this slide]")
            lines.append("")

    elif ext in (".png", ".jpg", ".jpeg", ".webp"):
        console.print(f"[dim]  Image file — sending to AI vision.[/dim]")
        with open(file_path, "rb") as f:
            raw = f.read()
        img_bytes = resize_image(raw, max_dim)
        mime = "image/png" if ext == ".png" else "image/jpeg"
        image_parts.append((img_bytes, mime, Path(file_path).name))
        lines.append(f"## Image: {Path(file_path).name}")
        lines.append("[Image slide — attached for AI visual analysis]")

    else:
        lines.append(f"[Unsupported file type: {ext}]")

    return "\n".join(lines), image_parts

# --- AI ANALYSIS ---
def analyze_slides(content_text: str, image_parts: list, prompt_file: str, config: dict, status=None, target_language: str = "auto"):
    provider = config["SUMMARIZER_PROVIDER"]

    try:
        with open(prompt_file, "r", encoding="utf-8") as f:
            instructions = f.read()
    except FileNotFoundError:
        instructions = "Please analyze these lecture slides and provide a structured summary."

    if target_language == "english":
        instructions += "\n\nIMPORTANT: Please generate the final output in English."
    elif target_language == "malay":
        instructions += "\n\nIMPORTANT: Sila hasilkan output dalam Bahasa Melayu."

    has_images = bool(image_parts)

    if provider == "openai":
        client = OpenAI(api_key=config["OPENAI_API_KEY"])
        models = list(config.get("OPENAI_MODELS", ["gpt-4o-mini"]))
        if has_images and not any("gpt-4o" in m for m in models):
            models = ["gpt-4o"] + models
        last_error = ""
        for model_name in models:
            if status:
                status.update(f"[bold gold3]Analyzing with OpenAI ({model_name})...")
            try:
                if has_images and ("gpt-4o" in model_name or "vision" in model_name):
                    parts = [{"type": "text", "text": f"{instructions}\n\nSlide Content:\n{content_text}"}]
                    for img_bytes, mime_type, label in image_parts:
                        b64 = base64.b64encode(img_bytes).decode()
                        parts.append({"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{b64}"}})
                    messages = [{"role": "user", "content": parts}]
                else:
                    messages = [
                        {"role": "system", "content": instructions},
                        {"role": "user", "content": f"Analyze the following slide content:\n\n{content_text}"}
                    ]
                resp = client.chat.completions.create(model=model_name, messages=messages)
                return resp.choices[0].message.content
            except Exception as e:
                last_error = str(e)
                if status:
                    console.print(f"[gold3]⚠ OpenAI {model_name} failed: {last_error}. Trying next...[/gold3]")
        raise Exception(f"All OpenAI models failed. Last error: {last_error}")

    elif provider == "gemini":
        models = list(config.get("GEMINI_MODELS", ["gemini-2.0-flash"]))
        last_error = ""
        for model_name in models:
            if status:
                status.update(f"[bold gold3]Analyzing with Gemini ({model_name})...")
            try:
                client = genai.Client(api_key=config["GEMINI_API_KEY"])
                parts = [genai_types.Part.from_text(text=f"{instructions}\n\nSlide Content:\n{content_text}")]
                for img_bytes, mime_type, label in image_parts:
                    parts.append(genai_types.Part.from_bytes(data=img_bytes, mime_type=mime_type))
                response = client.models.generate_content(
                    model=model_name,
                    contents=genai_types.Content(parts=parts, role="user")
                )
                return response.text
            except Exception as e:
                last_error = str(e)
                if status:
                    console.print(f"[gold3]⚠ Gemini {model_name} failed: {last_error}. Trying next...[/gold3]")
        raise Exception(f"All Gemini models failed. Last error: {last_error}")

    raise Exception(f"Unsupported provider: {provider}")

# --- FILE / PROMPT HELPERS ---
def get_files(config: dict):
    files = []
    for pattern in config["FILE_PATTERN"].split(','):
        files.extend(glob.glob(os.path.join(config["INPUT_DIR"], pattern.strip())))
    files.sort(key=os.path.getmtime, reverse=True)
    return files

def get_prompts(config: dict):
    prompts = glob.glob(os.path.join(config["PROMPTS_DIR"], "*.md"))
    prompts.sort()
    return prompts

def clean_path(path_str: str) -> str:
    """Robustly clean a path string from macOS drag-and-drop input."""
    path_str = path_str.strip()
    # Remove leading '@' (some terminals prefix with this)
    if path_str.startswith('@'):
        path_str = path_str[1:]
    # Strip surrounding quotes (single or double)
    path_str = path_str.strip("'\"")
    # Strip again in case there were nested quotes
    path_str = path_str.strip()
    # Expand home directory and resolve to absolute path
    path_str = os.path.expanduser(path_str)
    # Unescape common shell escape sequences (e.g. backslash-space)
    for escaped, real in [('\\ ', ' '), ('\\(', '('), ('\\)', ')'), ("\\'" , "'")]:
        path_str = path_str.replace(escaped, real)
    return os.path.abspath(path_str)

def parse_multiple_paths(input_str: str):
    input_str = input_str.strip()
    if not input_str:
        return []
    # Insert space between adjacent quoted paths (e.g. 'path1''path2')
    input_str = re.sub(r"' *'", "' '", input_str)
    input_str = re.sub(r'" *"', '" "', input_str)
    try:
        paths = shlex.split(input_str)
    except ValueError:
        # Fallback: brute-force strip quotes and split
        paths = [input_str.strip("'\"")]
    return [clean_path(p) for p in paths if p.strip()]

def get_files_from_path(path: str, config: dict, strict_pattern: bool = True):
    """Resolve a path to a list of slide files.
    In strict_pattern mode (directory scans), only FILE_PATTERN types are returned.
    When strict_pattern=False (manual drop), any existing file is accepted.
    """
    if not os.path.exists(path):
        return []
    if os.path.isfile(path):
        if not strict_pattern:
            return [path]  # Accept any file the user explicitly dropped
        # Still accept if extension matches FILE_PATTERN
        ext = Path(path).suffix.lower()
        allowed = {p.strip().lstrip('*').lower() for p in config["FILE_PATTERN"].split(',')}
        if ext in allowed:
            return [path]
        return [path]  # Accept anyway — user explicitly chose this file
    if os.path.isdir(path):
        files = []
        for pattern in config["FILE_PATTERN"].split(','):
            # Case-insensitive glob: match both lower and upper extensions
            files.extend(glob.glob(os.path.join(path, pattern.strip())))
            files.extend(glob.glob(os.path.join(path, pattern.strip().upper())))
        files = list(dict.fromkeys(files))  # deduplicate
        files.sort()
        return files
    return []

# --- SUMMARY REPORT ---
def show_summary_report(times: dict, file_name: str, out_dir: str, prompt_name: str, provider: str):
    console.print("\n")
    console.print(Panel(f"[bold spring_green3]Processing Complete for[/bold spring_green3] [magenta3]{file_name}[/magenta3]", border_style="spring_green3"))
    table = Table(show_header=False, box=None)
    table.add_column("Key", style="bold magenta1")
    table.add_column("Value", style="white")
    table.add_row("Output Location", f"[gold3]{out_dir}[/gold3]")
    table.add_row("Extraction Time", f"{times.get('extract', 0):.2f} seconds")
    table.add_row("AI Analysis Time", f"{times.get('analyze', 0):.2f} seconds")
    table.add_row("Prompt Used", prompt_name)
    table.add_row("AI Provider", provider)
    console.print(table)
    console.print("\n")

# --- PROCESS FILE ---
def process_file(file_path: str, prompt_file: str, config: dict, target_language: str = "auto"):
    base_name = Path(file_path).stem
    out_dir = os.path.join(config["OUTPUT_DIR"], base_name)
    os.makedirs(out_dir, exist_ok=True)
    times = {}

    console.print(f"\n[bold magenta1]▶ Processing:[/bold magenta1] {Path(file_path).name}")
    console.print("[bold magenta1]▶ Extracting slide content...[/bold magenta1]")

    start = time.time()
    with Status("[bold gold3]Reading slides...", console=console) as status:
        try:
            content_text, image_parts = extract_slides_content(file_path, config)
            status.update("[bold spring_green3]✓ Content extracted!")
        except Exception as e:
            console.print(f"[bold red]✗ Extraction failed:[/bold red] {e}")
            console.print(f"[dim]{traceback.format_exc()}[/dim]")
            return
    times['extract'] = time.time() - start

    content_path = os.path.join(out_dir, "slides_content.md")
    with open(content_path, "w", encoding="utf-8") as f:
        f.write(f"# Extracted Slide Content: {Path(file_path).name}\n\n")
        f.write(content_text)
    console.print(f"[bold spring_green3]✓ Slide content saved:[/bold spring_green3] {content_path}")

    prompt_name = Path(prompt_file).stem
    console.print(f"\n[bold magenta1]▶ AI Processing with {config['SUMMARIZER_PROVIDER'].title()}...[/bold magenta1]")
    if image_parts:
            console.print(f"[dim]  Sending {len(image_parts)} image page(s) for visual analysis.[/dim]")
    with Status(f"[bold gold3]Applying prompt '{prompt_name}'...", console=console) as status:
        start = time.time()
        try:
            analysis = analyze_slides(content_text, image_parts, prompt_file, config, status=status, target_language=target_language)
            times['analyze'] = time.time() - start
            analysis_path = os.path.join(out_dir, f"{prompt_name}.md")
            with open(analysis_path, "w", encoding="utf-8") as f:
                f.write(analysis)
            status.update(f"[bold spring_green3]✓ Saved: {analysis_path}")
            console.print("[bold spring_green3]✓ AI Analysis complete.[/bold spring_green3]")
        except Exception as e:
            times['analyze'] = time.time() - start
            console.print(f"[bold red]✗ AI Analysis failed:[/bold red] {e}")

    show_summary_report(times, Path(file_path).name, out_dir, prompt_name, config['SUMMARIZER_PROVIDER'])

# --- DASHBOARD ---
def dashboard():
    console.print(Panel(
        "[bold magenta1]SlidSum CLI[/bold magenta1]\n"
        "[dim]Lecture Slides Content Extractor & AI Analyzer[/dim]\n"
        "[dim magenta3]by kahmeng  kahmeng15.github.io[/dim magenta3]",
        border_style="magenta3"
    ))
    config = setup_environment()

    while True:
        # Step 1: Mode
        mode_sel = inquirer.prompt([
            inquirer.List('mode', message="1. Select Processing Mode", choices=[
                ("Single File (Process one slide file)", "single"),
                ("Batch Processing (Process multiple files or a folder)", "batch"),
                ("Exit", "exit"),
            ])
        ])
        if not mode_sel or mode_sel['mode'] == 'exit':
            break
        global_mode = mode_sel['mode']

        files = get_files(config)
        prompts = get_prompts(config)
        target_files = []

        if not prompts:
            console.print("[bold red]Error:[/bold red] No prompt files found in the prompts/ folder. Add at least one .md file.")
            continue

        # Step 2: File selection
        if global_mode == "single":
            choices = []
            if files:
                choices.append(("Latest File", files[0]))
                if Separator: choices.append(Separator())
            choices.append(("Manual Path / Drag-and-Drop", "manual"))
            if files:
                if Separator: choices.append(Separator())
                for f in files:
                    mtime = time.strftime('%Y-%m-%d %H:%M', time.localtime(os.path.getmtime(f)))
                    choices.append((f"{Path(f).name} ({mtime})", f))

            ans = inquirer.prompt([inquirer.List('target_file', message="2. Select a file to process", choices=choices)])
            if not ans: continue

            if ans['target_file'] == "manual":
                path_input = input("Drag and drop your slide file here: ").strip()
                if not path_input: continue
                parsed = parse_multiple_paths(path_input)
                for p in parsed:
                    found = get_files_from_path(p, config, strict_pattern=False)
                    if found:
                        target_files = [found[0]]
                        break
                if not target_files:
                    attempted = parsed[0] if parsed else path_input
                    console.print(f"[bold red]Error:[/bold red] No valid file found.")
                    console.print(f"[dim]  Resolved path tried: {attempted}[/dim]")
                    console.print(f"[dim]  Tip: Make sure the file exists and the path has no trailing characters.[/dim]")
                    continue
            else:
                target_files = [ans['target_file']]

        elif global_mode == "batch":
            batch_choices = [("Drag and Drop Files/Folders (Collection Mode)", "folder_manual")]
            if files:
                if Separator: batch_choices.append(Separator())
                batch_choices.append(("Select Multiple Files from input_slides", "select_files"))

            batch_ans = inquirer.prompt([inquirer.List('type', message="2. Batch Input Method", choices=batch_choices)])
            if not batch_ans: continue

            if batch_ans['type'] == "folder_manual":
                console.print("\n[bold magenta1]▶ Collection Mode Enabled[/bold magenta1]")
                console.print("[dim]Drag and drop files/folders, press Enter after each.[/dim]")
                console.print("[dim]Type [bold gold3]'d'[/bold gold3] and press Enter when done.[/dim]\n")
                while True:
                    val = input(f"Add to queue ({len(target_files)} collected, 'd' to finish): ").strip()
                    if val.lower() == 'd': break
                    if not val: continue
                    added = 0
                    for p in parse_multiple_paths(val):
                        found = get_files_from_path(p, config, strict_pattern=False)
                        if found:
                            target_files.extend(found)
                            added += len(found)
                        else:
                            console.print(f"[gold3]⚠ No valid files found at: {p}[/gold3]")
                    if added > 0:
                        console.print(f"[bold spring_green3]✓ Added {added} file(s).[/bold spring_green3]")
                target_files = list(dict.fromkeys(target_files))
                if not target_files:
                    console.print(f"[bold gold3]Queue is empty. Returning to menu...[/bold gold3]")
                    continue

            elif batch_ans['type'] == "select_files":
                cb_choices = []
                for f in files:
                    mtime = time.strftime('%Y-%m-%d %H:%M', time.localtime(os.path.getmtime(f)))
                    cb_choices.append((f"{Path(f).name} ({mtime})", f))
                sel = inquirer.prompt([inquirer.Checkbox('files', message="Select files (Space to select, Enter to confirm)", choices=cb_choices)])
                if not sel or not sel['files']: continue
                target_files = sel['files']

        # Step 3: Prompt selection (mandatory — no skip option)
        prompt_choices = []
        for p in prompts:
            prompt_choices.append((Path(p).name, p))

        prompt_ans = inquirer.prompt([
            inquirer.List('target_prompt',
                message=f"3. Select an AI prompt (for {len(target_files)} file{'s' if len(target_files) > 1 else ''})",
                choices=prompt_choices)
        ])
        if not prompt_ans: break
        target_prompt = prompt_ans['target_prompt']

        # Step 4: Output language
        lang_ans = inquirer.prompt([
            inquirer.List('lang', message="4. Select AI output language", choices=[
                ("Same as Slides (Auto)", "auto"),
                ("English", "english"),
                ("Malay (Bahasa Melayu)", "malay"),
            ])
        ])
        if not lang_ans: break
        target_language = lang_ans['lang']

        # Process queue
        for i, file_path in enumerate(target_files):
            if len(target_files) > 1:
                console.print(Panel(
                    f"[bold magenta1]Queue: {i+1}/{len(target_files)}[/bold magenta1]\n[dim]{Path(file_path).name}[/dim]",
                    border_style="magenta3"
                ))

            # Overwrite check (single file only)
            base_name = Path(file_path).stem
            out_dir = os.path.join(config["OUTPUT_DIR"], base_name)
            existing = os.path.join(out_dir, "slides_content.md")
            if os.path.exists(existing) and len(target_files) == 1:
                console.print(f"\n[bold gold3]⚠ Output for '{base_name}' already exists![/bold gold3]")
                ow_ans = inquirer.prompt([inquirer.Confirm('overwrite', message="Overwrite it?", default=False)])
                if not ow_ans or not ow_ans['overwrite']:
                    console.print("[dim]Skipping...[/dim]\n")
                    continue

            try:
                process_file(file_path, target_prompt, config, target_language=target_language)
            except Exception as e:
                console.print(Panel(
                    f"[bold red]Error processing '{Path(file_path).name}':[/bold red]\n\n{e}\n\n[dim]{traceback.format_exc()}[/dim]",
                    border_style="red"
                ))

        console.print(f"\n[bold spring_green3]✓ Done! Processed {len(target_files)} file{'s' if len(target_files) > 1 else ''}.[/bold spring_green3]")

        post = inquirer.prompt([
            inquirer.List('action', message="What next?", choices=[
                ("Return to Dashboard", "dash"),
                ("Exit", "exit"),
            ])
        ])
        if not post or post['action'] == 'exit':
            break

if __name__ == "__main__":
    try:
        dashboard()
    except KeyboardInterrupt:
        console.print("\n[bold gold3]Exiting...[/bold gold3]")
        sys.exit(0)
